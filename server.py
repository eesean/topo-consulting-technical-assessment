from flask import Flask, jsonify, render_template
import pandas as pd
import pymupdf
from pptx import Presentation
import json
import os
import matplotlib.pyplot as plt
from io import BytesIO
import base64
from flask_cors import CORS

app = Flask(__name__)
CORS(app)

class DataIngestion:
    def read_csv(file_path):
        return pd.read_csv(file_path)

    def read_json(file_path):
        with open(file_path, 'r') as file:
            data = json.load(file)
        return pd.json_normalize(data)

    def read_pdf(file_path):
        doc = pymupdf.open(file_path)
        text = []
        for page in doc:
            text.append(page.get_text())
        return pd.DataFrame({'Content': text})

    def read_pptx(file_path):
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return pd.DataFrame({'Content': text})

class DataProcessor:
    def clean_data(df):
        return df.dropna().reset_index(drop=True)

    def merge_data(dfs):
        return pd.concat(dfs, ignore_index=True)

class Visualization:
    def generate_bar_chart(data, column):
        plt.figure(figsize=(10, 6))
        data[column].value_counts().plot(kind='bar')
        plt.title(f"Bar Chart for {column}")
        plt.xlabel(column)
        plt.ylabel("Frequency")
        buffer = BytesIO()
        plt.savefig(buffer, format="png")
        buffer.seek(0)
        encoded_image = base64.b64encode(buffer.read()).decode('utf-8')
        buffer.close()
        return encoded_image

@app.route('/')
def index():
    csv_data = DataIngestion.read_csv("data.csv")
    cleaned_csv = DataProcessor.clean_data(csv_data)
    chart = Visualization.generate_bar_chart(cleaned_csv, csv_data.columns[0])
    return render_template("index.html", chart=chart)

@app.route('/api/data', methods=['GET'])
def get_data():
    csv_data = DataIngestion.read_csv("data.csv")
    json_data = DataIngestion.read_json("data.json")
    pdf_data = DataIngestion.read_pdf("data.pdf")
    pptx_data = DataIngestion.read_pptx("data.pptx")

    merged = DataProcessor.merge_data([csv_data, json_data, pdf_data, pptx_data])
    cleaned = DataProcessor.clean_data(merged)
    return jsonify(cleaned.to_dict(orient='records'))

@app.route('/api/data/<file_type>', methods=['GET'])
def get_data_by_type(file_type):
    if file_type == 'csv':
        data = DataIngestion.read_csv("data.csv")
    elif file_type == 'json':
        data = DataIngestion.read_json("data.json")
    elif file_type == 'pdf':
        data = DataIngestion.read_pdf("data.pdf")
    elif file_type == 'pptx':
        data = DataIngestion.read_pptx("data.pptx")
    else:
        return jsonify({"error": "Unsupported file type"}), 400
    cleaned = DataProcessor.clean_data(data)
    return jsonify(cleaned.to_dict(orient='records'))

if __name__ == '__main__':
    app.run(debug=True)